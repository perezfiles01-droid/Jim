# -*- coding: utf-8 -*-
"""Prove the deck changed only in text: same slides, same shapes, same geometry."""
from pptx import Presentation
SRC='/root/.claude/uploads/dff1b875-2cc7-5f42-9522-30eacbb8d668/66241d5d-EDRMS_Util_Report_2.pptx'
NEW='/home/user/Jim/EDRMS_Util_Report_3.pptx'

def profile(path):
    prs=Presentation(path); out=[]
    for i,s in enumerate(prs.slides,1):
        for sh in s.shapes:
            out.append((i, sh.shape_id, sh.name, sh.shape_type, sh.left, sh.top, sh.width, sh.height))
    return prs, out

def texts(path):
    prs=Presentation(path); out={}
    for i,s in enumerate(prs.slides,1):
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                out[(i,sh.shape_id,None,None)]=sh.text_frame.text
            if sh.has_table:
                for r,row in enumerate(sh.table.rows):
                    for c,cell in enumerate(row.cells):
                        out[(i,sh.shape_id,r,c)]=cell.text
    return out

a,pa=profile(SRC); b,pb=profile(NEW)
print("slides", len(a.slides.__iter__.__self__._sldIdLst), "->", len(b.slides.__iter__.__self__._sldIdLst))
print("slide size same:", (a.slide_width,a.slide_height)==(b.slide_width,b.slide_height))
print("shape count", len(pa), "->", len(pb))
geo=[x for x,y in zip(pa,pb) if x!=y]
print("GEOMETRY / SHAPE DIFFS:", geo if geo else "none")
pics_a=sum(1 for x in pa if str(x[3])=="PICTURE (13)"); pics_b=sum(1 for x in pb if str(x[3])=="PICTURE (13)")
print("pictures", pics_a, "->", pics_b)

ta,tb=texts(SRC),texts(NEW)
print("text nodes", len(ta), "->", len(tb), "| keys identical:", set(ta)==set(tb))
n=0
for k in ta:
    if ta[k]!=tb.get(k):
        n+=1
        loc=f"slide {k[0]}" + (f" table r{k[2]}c{k[3]}" if k[2] is not None else "")
        print(f"\n  [{loc}]\n    - {ta[k]}\n    + {tb.get(k)}")
print(f"\nTEXT CHANGES: {n}")
