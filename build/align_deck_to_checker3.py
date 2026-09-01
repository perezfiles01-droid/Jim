# -*- coding: utf-8 -*-
"""Align EDRMS_Util_Report to the two RAC sheets as revised in Checker_3.

Text only. Every edit is a substring swap inside an existing run so the run's
font, size, colour and the shape's geometry are untouched; nothing is added,
removed or moved. Run with --only N to apply one change set at a time.
"""
import sys, shutil, os
from pptx import Presentation

DECK = '/home/user/Jim/EDRMS_Util_Report_3.pptx'
SRC  = '/root/.claude/uploads/dff1b875-2cc7-5f42-9522-30eacbb8d668/66241d5d-EDRMS_Util_Report_2.pptx'

# (slide number, old, new). Slide numbers are 1 based.
CHANGES = {
 1: ("Slide 10 note: both sheets are assessed now", [
   (10,
    'This sheet has not yet been reviewed by RAC.  Bank-wide Oversight carries your '
    'comments against each row and they have been used above. Department Insights does '
    'not yet. The same review on this dashboard is the most useful next input we can receive.',
    'Both sheets now carry your assessment against every row and it has been used above.  '
    'What is still open is the decision itself: all 148 items across Bank-wide Oversight '
    'and Department Insights remain Pending Review. That is the most useful next input we can receive.'),
 ]),
 2: ("Department HQ, not department / office / RM", [
   (6,  'Every department, office and RM with its sites, documents, records and counterparts.',
        'Every department with its sites, documents, records and counterparts.'),
   (7,  'EDRMS sites, grouped by department, office or RM',
        'EDRMS sites, grouped by department HQ'),
   (8,  'The owning unit of each site',
        'The owning department. Narrowed to Department HQ on 28 August.'),
   (9,  "One unit's own view. Pick a department, office or RM and everything below refreshes for it.",
        "One department's own view. Pick a department and everything below refreshes for it."),
   (12, 'Two levels: department, then division. Countries, RMs and offices are divisions.',
        'Confirmed 28 August: Department HQ, then division. Countries, RMs and offices are divisions.'),
   (12, 'We read it as the department, office or RM on that row.',
        'We read it as the department on that row.'),
 ]),
 3: ("Physical counterpart, not paper counterpart", [
   (7,  'Declared records flagged as having a paper counterpart',
        'Declared records flagged as having a physical counterpart'),
   (10, 'Declared records flagged as having a paper counterpart',
        'Declared records flagged as having a physical counterpart'),
   (11, 'Declared records and paper counterparts per site',
        'Declared records and physical counterparts per site'),
 ]),
 4: ("The deleted window and category pickers", [
   (9, 'Page views and visited pages per site, over a window you choose.',
       'Page views and visited pages per site, over the reporting window.'),
   (9, "Libraries inside the unit's sites, grouped by file plan category.",
       "Libraries inside the department's sites, grouped by file plan category."),
 ]),
 5: ("Department, not unit, in the running prose", [
   (9,  'Seven tiles for the selected unit, each opening its own table.',
        'Seven tiles for the selected department, each opening its own table.'),
   (9,  'Every site the unit owns, with its owner, documents, records and counterparts.',
        'Every site the department owns, with its owner, documents, records and counterparts.'),
   (10, 'Sites provisioned to this unit',      'Sites provisioned to this department'),
   (10, 'People with recorded activity for this unit', 'People with recorded activity for this department'),
   (10, "File count across the unit's sites",  "File count across the department's sites"),
   (10, "Declared records for the unit's sites", "Declared records for the department's sites"),
   (10, 'When the unit went live on EDRMS',    'When the department went live on EDRMS'),
   (11, 'Each site the unit owns and who owns it', 'Each site the department owns and who owns it'),
 ]),
}


def frames(slide):
    """Every text frame on the slide, tables included."""
    for sh in slide.shapes:
        if sh.has_text_frame:
            yield sh.text_frame
        if sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    yield cell.text_frame


def swap(tf, old, new):
    """Replace old with new, keeping the run that carries it."""
    hits = 0
    for para in tf.paragraphs:
        if old not in para.text:
            continue
        for run in para.runs:                 # the common case: one run holds it
            if old in run.text:
                run.text = run.text.replace(old, new)
                hits += 1
                break
        else:                                 # split across runs: collapse onto the first
            first = para.runs[0]
            first.text = para.text.replace(old, new)
            for extra in para.runs[1:]:
                extra.text = ''
            hits += 1
    return hits


def apply(only=None):
    if not os.path.exists(DECK):
        shutil.copy(SRC, DECK)
    prs = Presentation(DECK)
    todo = CHANGES if only is None else {only: CHANGES[only]}
    for num, (label, edits) in sorted(todo.items()):
        print(f'\nChange {num}: {label}')
        for slide_no, old, new in edits:
            hits = sum(swap(tf, old, new) for tf in frames(prs.slides[slide_no - 1]))
            if hits != 1:
                raise SystemExit(f'  ABORT slide {slide_no}: {hits} matches for {old[:60]!r}')
            print(f'  slide {slide_no}  {old[:58]}...  ->  {new[:58]}...')
    prs.save(DECK)
    print(f'\nsaved {DECK}')


if __name__ == '__main__':
    only = int(sys.argv[sys.argv.index('--only') + 1]) if '--only' in sys.argv else None
    apply(only)
