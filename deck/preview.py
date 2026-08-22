"""Render each slide to HTML from the packed .pptx, using the real shape
geometry and text, so layout and overflow can be checked without LibreOffice
(soffice cannot load any file in this sandbox)."""
import sys, html, json
from pptx import Presentation
from pptx.util import Emu

SRC=sys.argv[1]; OUT=sys.argv[2]
prs=Presentation(SRC)
SW=prs.slide_width/914400.0; SH=prs.slide_height/914400.0
PX=96.0

def color_of(fmt):
    try:
        if fmt.type is not None and fmt.fore_color.type is not None:
            return '#'+str(fmt.fore_color.rgb)
    except Exception: pass
    return None

parts=[]
for idx,slide in enumerate(prs.slides,1):
    bg='#FFFFFF'
    try:
        c=slide.background.fill
        if c.type is not None: bg='#'+str(c.fore_color.rgb)
    except Exception: pass
    body=[]
    for sh in slide.shapes:
        try: x,y,w,h=sh.left/914400.0,sh.top/914400.0,sh.width/914400.0,sh.height/914400.0
        except Exception: continue
        st=f"left:{x*PX}px;top:{y*PX}px;width:{w*PX}px;height:{h*PX}px;"
        if sh.shape_type is not None and sh.has_text_frame is False or sh.shape_type==13:
            pass
        fill=None
        try: fill=color_of(sh.fill)
        except Exception: pass
        cls='shp'
        if sh.shape_type==13:  # picture
            import base64
            try:
                b=base64.b64encode(sh.image.blob).decode()
                body.append(f'<img class="pic" style="{st}" src="data:image/png;base64,{b}">')
            except Exception:
                body.append(f'<div class="pic" style="{st}"></div>')
            continue
        if fill: st+=f"background:{fill};"
        rad=''
        try:
            if 'roundRect' in str(sh.auto_shape_type): rad='border-radius:6px;'
            if 'OVAL' in str(sh.auto_shape_type) or 'ellipse' in str(sh.auto_shape_type): rad='border-radius:50%;'
        except Exception: pass
        st+=rad
        if not sh.has_text_frame:
            body.append(f'<div class="{cls}" style="{st}"></div>'); continue
        tf=sh.text_frame
        runs=[]
        for para in tf.paragraphs:
            for r in para.runs:
                f=r.font
                col='#333'
                try:
                    if f.color and f.color.type is not None: col='#'+str(f.color.rgb)
                except Exception: pass
                sz=f.size.pt if f.size else 12
                fam=f.name or 'Calibri'
                bold='font-weight:700;' if f.bold else ''
                ital='font-style:italic;' if f.italic else ''
                runs.append(f'<span style="font-size:{sz}pt;font-family:\'{fam}\',serif;color:{col};{bold}{ital}">{html.escape(r.text)}</span>')
            runs.append('<br>')
        if runs and runs[-1]=='<br>': runs.pop()
        al=''
        try:
            a=tf.paragraphs[0].alignment
            if a is not None and 'CENTER' in str(a): al='text-align:center;'
        except Exception: pass
        va=''
        try:
            v=str(tf.vertical_anchor)
            if 'MIDDLE' in v: va='display:flex;align-items:center;'
        except Exception: pass
        body.append(f'<div class="{cls} txt" style="{st}{al}{va}">{"".join(runs) or "&nbsp;"}</div>')
    parts.append(f'<div class="slide" style="background:{bg}" data-n="{idx}">'
                 f'<div class="num">{idx}</div>{"".join(body)}</div>')

open(OUT,'w').write(f"""<!doctype html><meta charset=utf-8><style>
body{{margin:0;background:#666;font-family:Calibri,Arial}}
.slide{{position:relative;width:{SW*PX}px;height:{SH*PX}px;margin:14px auto;overflow:visible;
 box-shadow:0 2px 10px rgba(0,0,0,.4)}}
.shp{{position:absolute;box-sizing:border-box;overflow:visible}}
.txt{{line-height:1.18;padding:0;overflow:hidden}}
.pic{{position:absolute;background:repeating-linear-gradient(45deg,#cfe0ee,#cfe0ee 8px,#bcd3e6 8px,#bcd3e6 16px);
 border:1px solid #9ab}}
.num{{position:absolute;right:-30px;top:0;color:#fff;font:700 13px Arial}}
</style>{''.join(parts)}""")
print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {OUT}")
